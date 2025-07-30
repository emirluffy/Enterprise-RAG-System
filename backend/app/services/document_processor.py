"""
Enterprise RAG System - Document Processing Service
Handles PDF, DOCX, TXT file processing and chunking for RAG
"""
import os
import uuid
import hashlib
from typing import List, Dict, Optional, BinaryIO
from pathlib import Path
import asyncio
from datetime import datetime

# Document processing libraries
try:
    import PyPDF2
    from docx import Document as DocxDocument
    from pptx import Presentation
    PDF_AVAILABLE = True
    DOCX_AVAILABLE = True
    PPTX_AVAILABLE = True
except ImportError as e:
    print(f"Document processing libraries not available: {e}")
    PDF_AVAILABLE = False
    DOCX_AVAILABLE = False
    PPTX_AVAILABLE = False

from app.services.embeddings import embeddings_service
from app.services.vector_store import vector_store_service

class DocumentProcessor:
    """Document processing service for Enterprise RAG"""
    
    def __init__(self, chunk_size: int = 300, chunk_overlap: int = 30):
        self.supported_formats = []
        if PDF_AVAILABLE:
            self.supported_formats.extend(['.pdf'])
        if DOCX_AVAILABLE:
            self.supported_formats.extend(['.docx'])
        if PPTX_AVAILABLE:
            self.supported_formats.extend(['.pptx', '.ppt'])
        self.supported_formats.extend(['.txt'])
        
        # Processing settings
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.max_file_size = 50 * 1024 * 1024  # 50MB limit (from PRD)
    
    def extract_text_from_pdf(self, file_content: bytes) -> str:
        """Extract text from PDF file"""
        if not PDF_AVAILABLE:
            raise Exception("PDF processing not available - PyPDF2 not installed")
        
        try:
            from io import BytesIO
            pdf_reader = PyPDF2.PdfReader(BytesIO(file_content))
            
            text = ""
            for page_num, page in enumerate(pdf_reader.pages):
                page_text = page.extract_text()
                if page_text.strip():
                    text += f"\n--- Page {page_num + 1} ---\n"
                    text += page_text + "\n"
            
            return text.strip()
        except Exception as e:
            raise Exception(f"Error extracting PDF text: {str(e)}")
    
    def extract_text_from_docx(self, file_content: bytes) -> str:
        """Extract text from DOCX file"""
        if not DOCX_AVAILABLE:
            raise Exception("DOCX processing not available - python-docx not installed")
        
        try:
            from io import BytesIO
            from zipfile import BadZipFile
            
            # First, check if it's a valid ZIP file (DOCX is compressed as ZIP)
            try:
                doc = DocxDocument(BytesIO(file_content))
            except BadZipFile:
                # More descriptive error for corrupted DOCX files
                raise Exception(
                    "Invalid DOCX file: File appears to be corrupted or not a valid DOCX format. "
                    "Please ensure the file is a genuine Microsoft Word document (.docx) and not renamed."
                )
            except Exception as zip_error:
                raise Exception(f"DOCX file format error: {str(zip_error)}")
            
            text = ""
            paragraph_count = 0
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    # Safe handling of special characters (Greek, Turkish, emoji, etc.)
                    try:
                        paragraph_text = paragraph.text.encode('utf-8', errors='replace').decode('utf-8')
                        text += paragraph_text + "\n"
                        paragraph_count += 1
                    except Exception as e:
                        print(f"⚠️ Character encoding warning in paragraph {paragraph_count}: {e}")
                        # Fallback: try to extract what we can
                        try:
                            safe_text = paragraph.text.encode('ascii', errors='ignore').decode('ascii')
                            if safe_text.strip():
                                text += safe_text + "\n"
                                paragraph_count += 1
                        except:
                            pass  # Skip problematic paragraphs
            
            # Also extract text from tables if any
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            try:
                                cell_text = cell.text.encode('utf-8', errors='replace').decode('utf-8')
                                text += cell_text + "\n"
                            except Exception as e:
                                print(f"⚠️ Character encoding warning in table cell: {e}")
                                # Try ASCII fallback for table cells too
                                try:
                                    safe_text = cell.text.encode('ascii', errors='ignore').decode('ascii')
                                    if safe_text.strip():
                                        text += safe_text + "\n"
                                except:
                                    pass
            
            print(f"✅ DOCX processed: {paragraph_count} paragraphs, {len(text)} characters")
            
            if not text.strip():
                raise Exception("DOCX file contains no readable text content")
            
            return text.strip()
            
        except Exception as e:
            # If it's already our custom exception, re-raise
            if "Invalid DOCX file" in str(e) or "DOCX file format error" in str(e):
                raise e
            # Otherwise, wrap in a more descriptive error
            raise Exception(f"Error extracting DOCX text: {str(e)}")
    
    def extract_text_from_txt(self, file_content: bytes) -> str:
        """Extract text from TXT file"""
        try:
            # Try UTF-8 first, fallback to latin-1
            try:
                return file_content.decode('utf-8')
            except UnicodeDecodeError:
                return file_content.decode('latin-1')
        except Exception as e:
            raise Exception(f"Error extracting TXT text: {str(e)}")
    
    def extract_text_from_pptx(self, file_content: bytes) -> str:
        """Extract text from PowerPoint (.pptx) file using Context7 verified patterns"""
        if not PPTX_AVAILABLE:
            raise Exception("PowerPoint processing not available - python-pptx not installed")
        
        try:
            from io import BytesIO
            
            # Load the presentation
            presentation = Presentation(BytesIO(file_content))
            
            # Context7 verified pattern for complete text extraction
            all_text_runs = []
            slide_count = 0
            
            for slide in presentation.slides:
                slide_count += 1
                slide_text_runs = []
                
                # Add slide marker for better chunking
                slide_text_runs.append(f"\n--- Slide {slide_count} ---")
                
                # Context7 Pattern: Complete text extraction from all shapes
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    
                    # Extract complete paragraphs to preserve context
                    shape_paragraphs = []
                    text_frame = getattr(shape, 'text_frame', None)
                    if not text_frame:
                        continue
                        
                    for paragraph in text_frame.paragraphs:
                        paragraph_runs = []
                        for run in paragraph.runs:
                            if run.text.strip():  # Only non-empty text
                                paragraph_runs.append(run.text)
                        
                        # Join runs within same paragraph (preserves full sentence context)
                        if paragraph_runs:
                            complete_paragraph = "".join(paragraph_runs).strip()
                            if complete_paragraph:
                                shape_paragraphs.append(complete_paragraph)
                    
                    # Join paragraphs with proper spacing
                    if shape_paragraphs:
                        shape_complete_text = "\n".join(shape_paragraphs)
                        slide_text_runs.append(shape_complete_text)
                
                # Also extract notes if available (Context7 pattern)
                try:
                    notes_slide = slide.notes_slide
                    if notes_slide and hasattr(notes_slide, 'notes_text_frame'):
                        notes_text_frame = getattr(notes_slide, 'notes_text_frame', None)
                        if notes_text_frame and hasattr(notes_text_frame, 'text'):
                            notes_text = notes_text_frame.text.strip()
                            if notes_text:
                                slide_text_runs.append(f"[Notes] {notes_text}")
                except Exception:
                    pass  # Skip notes if unavailable
                
                # Add all slide content to master collection
                if len(slide_text_runs) > 1:  # More than just slide marker
                    all_text_runs.extend(slide_text_runs)
            
            if not all_text_runs:
                print(f"⚠️ No text content found in PowerPoint file")
                return ""
            
            # Join all content with proper spacing to preserve context
            full_text = "\n\n".join(all_text_runs)
            
            print(f"✅ PowerPoint processed: {len(presentation.slides)} slides, {len(full_text)} characters")
            
            return full_text
            
        except Exception as e:
            raise Exception(f"Error extracting PowerPoint text: {str(e)}")
    
    def _extract_table_text_from_pptx(self, table) -> str:
        """Extract text from PowerPoint table"""
        try:
            table_content = []
            
            for row in table.rows:
                row_content = []
                for cell in row.cells:
                    cell_text = cell.text.strip() if cell.text else ""
                    row_content.append(cell_text)
                
                # Join cells with tab separator
                if any(row_content):  # Only add non-empty rows
                    table_content.append("\t".join(row_content))
            
            return "\n".join(table_content)
            
        except Exception as e:
            print(f"⚠️ Error extracting table content: {str(e)}")
            return ""
    
    def extract_text(self, file_content: bytes, filename: str) -> str:
        """Extract text from any supported file format"""
        file_ext = Path(filename).suffix.lower()
        
        if file_ext == '.pdf':
            return self.extract_text_from_pdf(file_content)
        elif file_ext == '.docx':
            return self.extract_text_from_docx(file_content)
        elif file_ext in ['.pptx', '.ppt']:
            return self.extract_text_from_pptx(file_content)
        elif file_ext == '.txt':
            return self.extract_text_from_txt(file_content)
        else:
            raise Exception(f"Unsupported file format: {file_ext}")
    
    def create_chunks(self, text: str, source: str) -> List[Dict]:
        """
        Create text chunks with metadata
        
        Args:
            text: Full document text
            source: Source filename
            
        Returns:
            List of chunks with metadata
        """
        try:
            print(f"📝 Creating chunks for {source}...")
            print(f"📊 Text length: {len(text)} characters")
            
            # Use embeddings service for chunking
            chunks = embeddings_service.chunk_text(
                text, 
                chunk_size=self.chunk_size, 
                overlap=self.chunk_overlap
            )
            
            print(f"✅ Created {len(chunks)} chunks")
            
            # Validate chunk sizes
            for i, chunk in enumerate(chunks):
                if len(chunk) > 2000:  # Warn if chunk is too large
                    print(f"⚠️ Warning: Chunk {i} is large ({len(chunk)} chars)")
            
            # Create unique timestamp for this upload session
            timestamp = str(uuid.uuid4())[:8]
            
            # Create chunks with metadata
            chunk_documents = []
            for i, chunk in enumerate(chunks):
                try:
                    # Create unique ID with timestamp to prevent duplicates
                    safe_chunk = chunk.encode('utf-8', errors='replace').decode('utf-8')
                    chunk_hash = hashlib.md5(safe_chunk.encode('utf-8', errors='ignore')).hexdigest()[:8]
                    chunk_id = f"{source}_{timestamp}_{i}_{chunk_hash}"
                    
                    chunk_documents.append({
                        "id": chunk_id,
                        "text": safe_chunk,
                        "metadata": {
                            "source": source,
                            "chunk_index": i,
                            "total_chunks": len(chunks),
                            "timestamp": timestamp,
                            "char_count": len(safe_chunk)
                        }
                    })
                except Exception as chunk_error:
                    print(f"⚠️ Error processing chunk {i}: {chunk_error}")
                    print(f"⚠️ Chunk content preview: {chunk[:100]}...")
                    continue
                    
            print(f"✅ Successfully created {len(chunk_documents)} chunk documents")
            return chunk_documents
            
        except Exception as e:
            print(f"❌ Error in create_chunks: {str(e)}")
            print(f"❌ Error type: {type(e).__name__}")
            print(f"📊 Text stats: {len(text)} chars, source: {source}")
            
            # More detailed error info
            import traceback
            print(f"❌ Full traceback:")
            traceback.print_exc()
            
            raise Exception(f"Chunk creation failed: {str(e)} ({type(e).__name__})")
    
    async def process_document(
        self, 
        file_content: bytes, 
        filename: str,
        metadata: Optional[Dict] = None
    ) -> Dict:
        """
        Complete document processing pipeline
        
        Args:
            file_content: Raw file bytes
            filename: Original filename
            metadata: Optional additional metadata
            
        Returns:
            Processing result with statistics
        """
        # Validate file size
        if len(file_content) > self.max_file_size:
            raise Exception(f"File too large: {len(file_content)} bytes > {self.max_file_size}")
        
        # Validate file format
        file_ext = Path(filename).suffix.lower()
        if file_ext not in self.supported_formats:
            raise Exception(f"Unsupported format: {file_ext}. Supported: {self.supported_formats}")
        
        try:
            # Step 1: Extract text
            print(f"Extracting text from {filename}...")
            text = self.extract_text(file_content, filename)
            
            if not text.strip():
                raise Exception("No text content found in document")
            
            # Step 2: Create chunks
            print(f"Creating chunks for {filename}...")
            chunks = self.create_chunks(text, filename)
            
            # Step 3: Generate embeddings
            print(f"Generating embeddings for {len(chunks)} chunks...")
            chunk_texts = [chunk["text"] for chunk in chunks]
            embeddings = await embeddings_service.create_embeddings(chunk_texts)
            
            # Step 4: Prepare documents for vector store
            documents_for_vector_store = []
            for chunk, embedding in zip(chunks, embeddings):
                documents_for_vector_store.append({
                    "id": chunk["id"],
                    "embedding": embedding,
                    "metadata": {
                        **chunk["metadata"],
                        **(metadata or {}),
                        "content": chunk["text"],
                        "upload_date": datetime.now().isoformat(),
                        "file_size": len(file_content),
                        "file_type": file_ext
                    }
                })
            
            # Step 5: Store in vector database
            print(f"Storing {len(documents_for_vector_store)} documents in vector database...")
            success = await vector_store_service.upsert_documents(documents_for_vector_store)
            
            return {
                "success": success,
                "filename": filename,
                "file_size": len(file_content),
                "text_length": len(text),
                "word_count": len(text.split()),
                "chunks_created": len(chunks),
                "embeddings_generated": len(embeddings),
                "vector_storage": "completed" if success else "failed",
                "processing_time": "calculated_elsewhere"  # TODO: Add timing
            }
            
        except Exception as e:
            return {
                "success": False,
                "filename": filename,
                "error": str(e),
                "chunks_created": 0,
                "embeddings_generated": 0
            }
    
    async def get_processing_status(self) -> Dict:
        """Get current processing status and statistics"""
        try:
            vector_stats = await vector_store_service.get_index_stats()
            return {
                "service": "document_processor",
                "supported_formats": self.supported_formats,
                "pdf_available": PDF_AVAILABLE,
                "docx_available": DOCX_AVAILABLE,
                "chunk_size": self.chunk_size,
                "chunk_overlap": self.chunk_overlap,
                "max_file_size": self.max_file_size,
                "vector_store_stats": vector_stats
            }
        except Exception as e:
            return {
                "service": "document_processor",
                "error": str(e),
                "supported_formats": self.supported_formats
            }

# Create global instance
document_processor = DocumentProcessor() 